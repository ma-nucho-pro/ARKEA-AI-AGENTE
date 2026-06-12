
import os
import subprocess, shutil
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import parse_xml
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.drawing.image import Image as XLImage

BASE_DIR = Path(os.getenv("ARKEA_BUNDLE_DIR", Path(__file__).resolve().parents[2]))
DATA_DIR = Path(os.getenv("ARKEA_DATA_DIR", "./data"))
DOCS_DIR = DATA_DIR / "generated" / "docs"
DOCS_DIR.mkdir(parents=True, exist_ok=True)

def _logo_path():
    candidates = [
        BASE_DIR / "frontend" / "assets" / "arkea-logo.png",
        BASE_DIR / "frontend" / "assets" / "icon-512.png",
        Path(__file__).resolve().parents[2] / "frontend" / "assets" / "arkea-logo.png",
    ]
    for p in candidates:
        try:
            if p.exists():
                return str(p)
        except Exception:
            pass
    return ""

def safe_name(title: str):
    return ''.join(c if c.isalnum() or c in (' ', '_', '-') else '_' for c in (title or "arkea"))[:50].replace(' ', '_') or 'arkea'

def _filename(title, ext, output_dir=None):
    folder = Path(output_dir) if output_dir else DOCS_DIR
    folder.mkdir(parents=True, exist_ok=True)
    return str(folder / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{safe_name(title)}.{ext}")

def _find_soffice():
    candidates = [
        shutil.which('soffice'),
        shutil.which('libreoffice'),
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
    ]
    for c in candidates:
        if c and Path(c).exists():
            return str(c)
    return ''

def convert_to_pdf_with_libreoffice(path: str):
    """Try to create a PDF preview next to the Office file using LibreOffice if installed."""
    try:
        src = Path(path)
        if not src.exists():
            return ''
        soffice = _find_soffice()
        if not soffice:
            return ''
        outdir = src.parent / '_preview_pdf'
        outdir.mkdir(parents=True, exist_ok=True)
        cmd = [soffice, '--headless', '--convert-to', 'pdf', '--outdir', str(outdir), str(src)]
        subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60)
        pdf = outdir / (src.stem + '.pdf')
        return str(pdf) if pdf.exists() else ''
    except Exception:
        return ''

def _set_cell_text(cell, text, color="111827", bold=False):
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Segoe UI"
            r.font.size = Pt(18)
            r.font.bold = bold
            r.font.color.rgb = RGBColor.from_string(color)

def create_docx(title: str, paragraphs: list[str], tables: list[list[list[str]]] | None = None, output_dir: str | None = None):
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = DocxInches(0.6)
    sec.bottom_margin = DocxInches(0.6)
    sec.left_margin = DocxInches(0.65)
    sec.right_margin = DocxInches(0.65)

    logo = _logo_path()
    header = sec.header
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if logo:
        try:
            hp.add_run().add_picture(logo, width=DocxInches(0.45))
        except Exception:
            pass
    hr = hp.add_run("  ARKEA AI")
    hr.bold = True
    hr.font.size = DocxPt(14)
    hr.font.color.rgb = DocxRGBColor(37, 99, 235)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run(title or "Documento ARKEA")
    run.bold = True
    run.font.size = DocxPt(24)
    run.font.color.rgb = DocxRGBColor(30, 58, 138)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = sub.add_run("Documento generado por ARKEA AI · editable en Microsoft Word")
    sr.font.size = DocxPt(10)
    sr.font.color.rgb = DocxRGBColor(100, 116, 139)

    for idx, p in enumerate(paragraphs or []):
        if idx == 0:
            h = doc.add_heading("Resumen", level=1)
            h.runs[0].font.color.rgb = DocxRGBColor(37, 99, 235)
        para = doc.add_paragraph(str(p))
        para.paragraph_format.space_after = DocxPt(8)
        para.paragraph_format.line_spacing = 1.15
        for r in para.runs:
            r.font.name = "Calibri"
            r.font.size = DocxPt(11)

    # Automatic branded structure
    for heading, body in [
        ("Objetivo", "Presentar la información solicitada de forma clara, visual y ordenada."),
        ("Desarrollo", "El contenido puede ampliarse, corregirse o adaptarse con nuevas instrucciones en ARKEA AI."),
        ("Recomendación", "Puedes pedirme agregar tablas, portada, referencias, gráficos, anexos o una versión académica.")
    ]:
        h = doc.add_heading(heading, level=1)
        h.runs[0].font.color.rgb = DocxRGBColor(6, 182, 212)
        para = doc.add_paragraph(body)
        for r in para.runs:
            r.font.name = "Calibri"
            r.font.size = DocxPt(11)

    if tables:
        for tbl in tables:
            if not tbl:
                continue
            t = doc.add_table(rows=len(tbl), cols=max(len(r) for r in tbl))
            t.alignment = WD_TABLE_ALIGNMENT.CENTER
            t.style = "Table Grid"
            for i, row in enumerate(tbl):
                for j, val in enumerate(row):
                    cell = t.cell(i, j)
                    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                    cell.text = str(val)
                    if i == 0:
                        for p in cell.paragraphs:
                            for r in p.runs:
                                r.bold = True
                                r.font.color.rgb = DocxRGBColor(255, 255, 255)
                        cell._tc.get_or_add_tcPr().append(parse_xml(r'<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" w:fill="2563EB"/>'))

    footer = sec.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer.add_run("ARKEA AI · robertmanuchojarapeche@gmail.com")
    fr.font.size = DocxPt(8)
    fr.font.color.rgb = DocxRGBColor(100, 116, 139)

    path = _filename(title, "docx", output_dir)
    doc.save(path)
    return path

def create_xlsx(title: str, rows: list[list[str]], sheets: list[dict] | None = None, output_dir: str | None = None):
    wb = Workbook()
    default = wb.active
    wb.remove(default)
    sheets_data = sheets or [{"name": title[:31] or "Hoja1", "rows": rows}]
    logo = _logo_path()

    for sh in sheets_data:
        ws = wb.create_sheet((sh.get("name") or "Hoja")[:31])
        data = sh.get("rows") or rows or [["Campo", "Valor"], ["Estado", "Generado por ARKEA AI"]]
        col_count = max(2, max(len(r) for r in data))
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=col_count)
        ws.cell(1,1).value = "ARKEA AI · " + (sh.get("name") or title)
        ws.cell(1,1).font = Font(name="Segoe UI", bold=True, color="FFFFFF", size=16)
        ws.cell(1,1).fill = PatternFill("solid", fgColor="1D4ED8")
        ws.cell(1,1).alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[1].height = 28
        if logo:
            try:
                img = XLImage(logo)
                img.width, img.height = 48, 48
                ws.add_image(img, "A2")
                start_row = 4
            except Exception:
                start_row = 3
        else:
            start_row = 3

        thin = Side(style="thin", color="94A3B8")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)
        for r_idx, row in enumerate(data, start=start_row):
            for c_idx in range(1, col_count+1):
                cell = ws.cell(r_idx, c_idx)
                cell.value = row[c_idx-1] if c_idx-1 < len(row) else ""
                cell.border = border
                cell.alignment = Alignment(vertical="center", wrap_text=True)
                if r_idx == start_row:
                    cell.font = Font(name="Segoe UI", bold=True, color="FFFFFF")
                    cell.fill = PatternFill("solid", fgColor="2563EB")
                    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                elif r_idx % 2 == 0:
                    cell.fill = PatternFill("solid", fgColor="F1F5F9")
        ws.freeze_panes = ws.cell(start_row+1, 1)
        ws.auto_filter.ref = f"A{start_row}:{ws.cell(start_row, col_count).coordinate}"
        for col in range(1, col_count+1):
            letter = ws.cell(1,col).column_letter
            max_len = 12
            for row in ws.iter_rows(min_col=col, max_col=col):
                for cell in row:
                    max_len = max(max_len, len(str(cell.value or "")))
            ws.column_dimensions[letter].width = min(45, max_len + 3)

    path = _filename(title, "xlsx", output_dir)
    wb.save(path)
    return path

def create_pptx(title: str, slides: list[dict], output_dir: str | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    logo = _logo_path()

    def bg(slide, color="07111F"):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color)

    def add_logo(slide):
        if logo:
            try:
                slide.shapes.add_picture(logo, Inches(0.35), Inches(0.28), width=Inches(0.48))
            except Exception:
                pass
        tx = slide.shapes.add_textbox(Inches(0.92), Inches(0.28), Inches(3), Inches(0.35))
        p = tx.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "ARKEA AI"
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(255,255,255)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, "08111F"); add_logo(s)
    title_box = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.4), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title or "Presentación ARKEA"
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = RGBColor(96,165,250)
    sub = s.shapes.add_textbox(Inches(0.95), Inches(3.05), Inches(10.8), Inches(0.8))
    pr = sub.text_frame.paragraphs[0].add_run()
    pr.text = "Generado por ARKEA AI · editable en PowerPoint"
    pr.font.size = Pt(22); pr.font.color.rgb = RGBColor(203,213,225)

    for item in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide, "0B1628"); add_logo(slide)
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.8), Inches(0.75))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = item.get("title", "Diapositiva")
        run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = RGBColor(125,211,252)
        body = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.1), Inches(4.6))
        tf = body.text_frame
        tf.word_wrap = True
        lines = str(item.get("body", "")).split("\n")
        for i, line in enumerate(lines[:8]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = 0
            r = p.add_run(); r.text = line.strip() or " "
            r.font.size = Pt(21); r.font.color.rgb = RGBColor(241,245,249)
        # accent bar
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(6.85), Inches(11.8), Inches(0.08))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(6,182,212)
        shape.line.fill.background()

    path = _filename(title, "pptx", output_dir)
    prs.save(path)
    return path
