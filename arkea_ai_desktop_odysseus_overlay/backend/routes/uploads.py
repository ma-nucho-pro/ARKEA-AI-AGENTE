
import os, json, shutil, base64, mimetypes
from pathlib import Path
from fastapi import APIRouter, UploadFile, File, Form
from backend.arkea_core.db import get_setting

router = APIRouter(prefix='/api/arkea/uploads', tags=['uploads'])

UPLOAD_ROOT = Path(os.getenv('ARKEA_DATA_DIR', './data')) / 'uploads'
UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)

IMAGE_EXTS = {'.png','.jpg','.jpeg','.webp','.gif','.bmp'}
VIDEO_EXTS = {'.mp4','.webm','.mov','.avi','.mkv'}
AUDIO_EXTS = {'.wav','.mp3','.m4a','.ogg','.webm'}
DOC_EXTS = {'.pdf','.docx','.doc','.xlsx','.xls','.pptx','.txt','.md','.csv','.json','.html','.js','.ts','.py','.zip'}

def safe_name(name: str):
    name = name or 'archivo'
    return ''.join(c if c.isalnum() or c in '._- ' else '_' for c in name)[:160]

def _read_text(path: Path, ext: str):
    try:
        if ext in {'.txt','.md','.csv','.json','.html','.css','.js','.ts','.py','.sql','.xml','.log','.ini'}:
            return path.read_text(encoding='utf-8', errors='ignore')[:220000]
        if ext == '.docx':
            from docx import Document
            d = Document(str(path))
            return '\n'.join(p.text for p in d.paragraphs if p.text.strip())[:220000]
        if ext in {'.xlsx','.xls'}:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets[:8]:
                out.append(f'--- Hoja: {ws.title} ---')
                for row in ws.iter_rows(max_row=250, values_only=True):
                    vals = [str(v) if v is not None else '' for v in row]
                    if any(vals):
                        out.append('\t'.join(vals))
            return '\n'.join(out)[:220000]
        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(str(path))
            out=[]
            for i, s in enumerate(prs.slides, start=1):
                out.append(f'--- Diapositiva {i} ---')
                for shape in s.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        out.append(shape.text.strip())
            return '\n'.join(out)[:220000]
        if ext == '.pdf':
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                out=[]
                for i, page in enumerate(reader.pages[:80], start=1):
                    out.append(f'--- Página {i} ---')
                    out.append(page.extract_text() or '')
                return '\n'.join(out)[:220000]
            except Exception as e:
                return f'[PDF subido, pero no pude extraer texto automáticamente: {e}]'
    except Exception as e:
        return f'[No pude extraer texto local: {e}]'
    return ''

def _data_url(path: Path, ext: str):
    mime = mimetypes.guess_type(str(path))[0] or ('image/png' if ext == '.png' else 'application/octet-stream')
    raw = path.read_bytes()
    if len(raw) > 8_000_000:
        return ''
    return f"data:{mime};base64,{base64.b64encode(raw).decode('ascii')}"

@router.post('/file')
async def upload_file(file: UploadFile = File(...), note: str = Form('')):
    raw = await file.read()
    fname = safe_name(file.filename or 'archivo')
    folder = UPLOAD_ROOT
    folder.mkdir(parents=True, exist_ok=True)
    path = folder / fname
    base = path.stem
    ext = path.suffix.lower()
    i = 1
    while path.exists():
        path = folder / f'{base}_{i}{ext}'
        i += 1
    path.write_bytes(raw)
    rel = '/data/uploads/' + path.name
    kind = 'file'
    preview_html = ''
    data_url = ''
    extracted_text = ''

    if ext in IMAGE_EXTS:
        kind = 'image'
        data_url = _data_url(path, ext)
        preview_html = f"""<!doctype html><html><body style='margin:0;background:#020617;color:white;font-family:Segoe UI,Arial'><img src='{rel}' style='max-width:100%;max-height:100vh;display:block;margin:auto'/><div style='padding:16px'>Imagen subida: {fname}</div></body></html>"""
    elif ext in VIDEO_EXTS:
        kind = 'video'
        preview_html = f"""<!doctype html><html><body style='margin:0;background:#020617;color:white;font-family:Segoe UI,Arial'><video src='{rel}' controls autoplay muted style='max-width:100%;max-height:92vh;display:block;margin:auto'></video><div style='padding:16px'>Video subido: {fname}</div></body></html>"""
    elif ext in AUDIO_EXTS:
        kind = 'audio'
        preview_html = f"""<!doctype html><html><body style='padding:24px;background:#020617;color:white;font-family:Segoe UI,Arial'><h1>Audio subido</h1><audio src='{rel}' controls></audio><p>{fname}</p></body></html>"""
    else:
        kind = 'document' if ext in DOC_EXTS else 'file'
        extracted_text = _read_text(path, ext)
        preview_html = f"""<!doctype html><html><body style='padding:24px;background:#020617;color:white;font-family:Segoe UI,Arial'><h1>Archivo subido</h1><p>{fname}</p><p>Ruta: {rel}</p><p>Texto extraído: {len(extracted_text)} caracteres.</p><pre style='white-space:pre-wrap;background:#111827;border-radius:14px;padding:14px'>{extracted_text[:2500]}</pre></body></html>"""

    return {
        'ok': True,
        'kind': kind,
        'filename': fname,
        'path': str(path),
        'url': rel,
        'data_url': data_url,
        'extracted_text': extracted_text,
        'html_content': preview_html,
        'message': f'Archivo subido: {fname}'
    }
